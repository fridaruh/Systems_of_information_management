"""
Clasificador de slides: determina si un slide es de solo texto o contiene contenido visual
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Tuple
from config import Config


class SlideClassifier:
    """Clasifica slides en 'text' o 'visual' para optimizar costos de API"""

    def __init__(self, config: Config = None):
        self.config = config or Config()

    def classify_presentation(self, pptx_path: str) -> Dict[int, str]:
        """
        Clasifica todos los slides de una presentación

        Args:
            pptx_path: Ruta al archivo PPTX

        Returns:
            Dict con {slide_number: 'text' | 'visual'}
        """
        prs = Presentation(pptx_path)
        classifications = {}

        for idx, slide in enumerate(prs.slides, start=1):
            classification = self._classify_slide(slide, idx)
            classifications[idx] = classification

            if self.config.VERBOSE:
                print(f"  Slide {idx}: {classification}")

        return classifications

    def _classify_slide(self, slide, slide_num: int) -> str:
        """
        Clasifica un slide individual

        Criterios para 'visual':
        1. Tiene imágenes o formas complejas
        2. El texto contiene keywords visuales
        3. Tiene tablas complejas
        """
        # Extraer información del slide
        has_images = self._has_images(slide)
        text_content = self._extract_text(slide)
        has_visual_keywords = self._contains_visual_keywords(text_content)
        has_tables = self._has_tables(slide)
        has_charts = self._has_charts(slide)

        # Lógica de clasificación
        if has_images or has_charts:
            return 'visual'

        if has_tables:
            return 'visual'

        if has_visual_keywords and len(text_content) < 100:
            # Si menciona contenido visual pero tiene poco texto, probablemente es visual
            return 'visual'

        if len(text_content) >= self.config.MIN_TEXT_LENGTH:
            return 'text'

        # Default: si no tiene texto suficiente ni elementos visuales claros
        return 'text' if text_content else 'visual'

    def _has_images(self, slide) -> bool:
        """Detecta si el slide tiene imágenes"""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            # Verificar si tiene fill de imagen
            if hasattr(shape, 'fill') and shape.fill.type == 6:  # PICTURE fill
                return True
        return False

    def _has_tables(self, slide) -> bool:
        """Detecta si el slide tiene tablas"""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return True
        return False

    def _has_charts(self, slide) -> bool:
        """Detecta si el slide tiene gráficos"""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return True
        return False

    def _extract_text(self, slide) -> str:
        """Extrae todo el texto de un slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
        return " ".join(text_parts).strip()

    def _contains_visual_keywords(self, text: str) -> bool:
        """Verifica si el texto contiene keywords que indican contenido visual"""
        text_lower = text.lower()
        return any(keyword in text_lower for keyword in self.config.VISUAL_KEYWORDS)

    def get_statistics(self, classifications: Dict[int, str]) -> Dict[str, int]:
        """Calcula estadísticas de clasificación"""
        stats = {
            'total': len(classifications),
            'text': sum(1 for c in classifications.values() if c == 'text'),
            'visual': sum(1 for c in classifications.values() if c == 'visual')
        }
        stats['text_percentage'] = (stats['text'] / stats['total'] * 100) if stats['total'] > 0 else 0
        return stats
