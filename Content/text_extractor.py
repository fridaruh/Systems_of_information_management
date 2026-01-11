"""
Extractor de texto de slides PPTX sin usar API (costo cero)
"""
from pptx import Presentation
from typing import Dict, List
from config import Config


class TextExtractor:
    """Extrae texto de slides clasificados como 'text-only'"""

    def __init__(self, config: Config = None):
        self.config = config or Config()

    def extract_from_presentation(
        self,
        pptx_path: str,
        slide_numbers: List[int]
    ) -> Dict[int, Dict[str, str]]:
        """
        Extrae texto de slides específicos

        Args:
            pptx_path: Ruta al archivo PPTX
            slide_numbers: Lista de números de slide a procesar

        Returns:
            Dict con {slide_num: {'title': str, 'content': str}}
        """
        prs = Presentation(pptx_path)
        extracted_data = {}

        for slide_num in slide_numbers:
            try:
                slide = prs.slides[slide_num - 1]  # Convert to 0-based index
                data = self._extract_slide_text(slide, slide_num)
                extracted_data[slide_num] = data

                if self.config.VERBOSE:
                    print(f"  Extraído slide {slide_num} (texto)")

            except IndexError:
                print(f"  Advertencia: Slide {slide_num} no existe")
                continue

        return extracted_data

    def _extract_slide_text(self, slide, slide_num: int) -> Dict[str, str]:
        """
        Extrae el contenido de texto de un slide individual

        Returns:
            Dict con 'title' y 'content'
        """
        title = ""
        content_parts = []

        for shape in slide.shapes:
            # Intentar extraer el título del slide
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()

                # El primer shape con texto suele ser el título
                if not title and len(text) < 150:
                    title = text
                else:
                    content_parts.append(text)

        # Si no encontramos título, usar el primero como título
        if not title and content_parts:
            title = content_parts.pop(0)

        # Procesar bullets y listas
        content = self._format_content(content_parts)

        return {
            'title': title or f"Slide {slide_num}",
            'content': content
        }

    def _format_content(self, content_parts: List[str]) -> str:
        """
        Formatea el contenido para markdown

        Convierte bullets y estructura de PowerPoint a markdown
        """
        formatted = []

        for part in content_parts:
            lines = part.split('\n')
            formatted_lines = []

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Detectar bullets (PowerPoint usa caracteres especiales)
                if line.startswith(('•', '-', '*', '·')):
                    line = f"- {line[1:].strip()}"
                elif self._looks_like_bullet(line):
                    line = f"- {line}"

                formatted_lines.append(line)

            if formatted_lines:
                formatted.append('\n'.join(formatted_lines))

        return '\n\n'.join(formatted)

    def _looks_like_bullet(self, line: str) -> bool:
        """
        Detecta si una línea parece ser un bullet point sin marcador
        """
        # Líneas cortas sin puntuación final suelen ser bullets
        if len(line) < 100 and not line.endswith(('.', '?', '!')):
            # No empieza con mayúscula seguida de más de 3 palabras
            words = line.split()
            if len(words) <= 10:
                return True
        return False

    def extract_speaker_notes(self, pptx_path: str, slide_num: int) -> str:
        """
        Extrae las notas del presentador si existen

        Args:
            pptx_path: Ruta al archivo PPTX
            slide_num: Número de slide

        Returns:
            Notas del presentador o string vacío
        """
        try:
            prs = Presentation(pptx_path)
            slide = prs.slides[slide_num - 1]

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                return text_frame.text.strip()

        except Exception as e:
            if self.config.VERBOSE:
                print(f"  No se pudieron extraer notas del slide {slide_num}: {e}")

        return ""
