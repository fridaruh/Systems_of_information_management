"""
Procesador de slides visuales usando Claude API
"""
import io
import base64
from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from PIL import Image
import anthropic
from config import Config


class VisualProcessor:
    """Procesa slides con contenido visual usando Claude Vision"""

    def __init__(self, config: Config = None):
        self.config = config or Config()
        self.client = anthropic.Anthropic(api_key=self.config.ANTHROPIC_API_KEY)
        self.total_input_tokens = 0
        self.total_output_tokens = 0

    def process_visual_slides(
        self,
        pptx_path: str,
        slide_numbers: List[int]
    ) -> Dict[int, Dict[str, str]]:
        """
        Procesa slides visuales enviándolos a Claude

        Args:
            pptx_path: Ruta al archivo PPTX
            slide_numbers: Lista de números de slides a procesar

        Returns:
            Dict con {slide_num: {'title': str, 'content': str, 'visual_description': str}}
        """
        prs = Presentation(pptx_path)
        processed_data = {}

        for slide_num in slide_numbers:
            try:
                slide = prs.slides[slide_num - 1]
                data = self._process_single_slide(slide, slide_num, pptx_path)
                processed_data[slide_num] = data

                if self.config.VERBOSE:
                    print(f"  Procesado slide {slide_num} (visual) - "
                          f"Tokens: {data.get('tokens_used', 0)}")

            except IndexError:
                print(f"  Advertencia: Slide {slide_num} no existe")
                continue
            except Exception as e:
                print(f"  Error procesando slide {slide_num}: {e}")
                # Fallback: extraer solo texto
                processed_data[slide_num] = {
                    'title': f'Slide {slide_num}',
                    'content': self._extract_text_fallback(slide),
                    'visual_description': 'Error al procesar contenido visual',
                    'error': str(e)
                }

        return processed_data

    def _process_single_slide(
        self,
        slide,
        slide_num: int,
        pptx_path: str
    ) -> Dict[str, str]:
        """
        Procesa un slide visual individual con Claude

        Returns:
            Dict con title, content, visual_description
        """
        # Convertir slide a imagen
        image_data = self._slide_to_image(slide, slide_num, pptx_path)

        # Extraer texto básico del slide
        slide_text = self._extract_text_fallback(slide)

        # Preparar prompt para Claude
        prompt = self._build_analysis_prompt(slide_text, slide_num)

        # Enviar a Claude
        response = self.client.messages.create(
            model=self.config.ANTHROPIC_MODEL,
            max_tokens=2000,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": image_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ],
                }
            ],
        )

        # Tracking de tokens
        self.total_input_tokens += response.usage.input_tokens
        self.total_output_tokens += response.usage.output_tokens

        # Parsear respuesta
        analysis = response.content[0].text

        return {
            'title': self._extract_title_from_text(slide_text) or f'Slide {slide_num}',
            'content': slide_text,
            'visual_description': analysis,
            'tokens_used': response.usage.input_tokens + response.usage.output_tokens
        }

    def _slide_to_image(self, slide, slide_num: int, pptx_path: str) -> str:
        """
        Convierte un slide a imagen PNG en base64

        Nota: python-pptx no soporta renderizado directo.
        Esta es una implementación simplificada que extrae imágenes embebidas.
        Para renderizado real, se necesitaría LibreOffice o similar.
        """
        # TODO: Implementar renderizado real de slides
        # Por ahora, retornamos un placeholder o usamos exportación manual

        # Opción 1: Extraer la primera imagen del slide si existe
        for shape in slide.shapes:
            if shape.shape_type == 6:  # PICTURE
                image = shape.image
                image_bytes = image.blob

                # Convertir a PNG si es necesario
                img = Image.open(io.BytesIO(image_bytes))
                buffer = io.BytesIO()
                img.save(buffer, format='PNG')
                image_data = base64.b64encode(buffer.getvalue()).decode('utf-8')

                return image_data

        # Opción 2: Crear imagen placeholder (para desarrollo)
        # En producción, usar herramienta externa para renderizar
        placeholder = Image.new('RGB', (960, 540), color='white')
        buffer = io.BytesIO()
        placeholder.save(buffer, format='PNG')
        return base64.b64encode(buffer.getvalue()).decode('utf-8')

    def _build_analysis_prompt(self, slide_text: str, slide_num: int) -> str:
        """Construye el prompt para análisis de Claude"""
        return f"""You are analyzing educational slide content for a knowledge base about Information Systems.

Slide #{slide_num}
Extracted text: {slide_text if slide_text else "(no text)"}

TASK: Extract and explain ONLY the educational concepts, information, and knowledge presented in this slide.

FOCUS ON:
- Core concepts and definitions
- Technical information and explanations
- Relationships between concepts
- Examples and code snippets (preserve exactly with proper syntax highlighting)
- Tables and structured data (convert to clean markdown tables)
- Diagram content (explain what concepts/relationships they illustrate, not their appearance)
- Practical applications and use cases

DO NOT INCLUDE:
- Descriptions of colors, fonts, backgrounds, or layout
- Comments about slide aesthetics or design
- Phrases like "the slide shows", "the diagram displays", "visually"
- Any metadata about how information is presented

FORMAT:
- Write in clear, concise English
- Use markdown formatting (headers, lists, code blocks, tables)
- Focus on knowledge useful for understanding the topic
- Be direct and informative
- Structure logically with appropriate headers

Output the educational content directly without preambles."""

    def _extract_text_fallback(self, slide) -> str:
        """Extrae texto básico de un slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())
        return "\n\n".join(text_parts)

    def _extract_title_from_text(self, text: str) -> str:
        """Extrae el título del texto del slide"""
        if not text:
            return ""
        lines = text.split('\n')
        # El primer línea no vacía suele ser el título
        for line in lines:
            line = line.strip()
            if line and len(line) < 150:
                return line
        return ""

    def get_cost_estimate(self) -> Dict[str, float]:
        """Calcula el costo estimado del procesamiento"""
        input_cost = self.total_input_tokens * self.config.COST_PER_INPUT_TOKEN
        output_cost = self.total_output_tokens * self.config.COST_PER_OUTPUT_TOKEN
        total_cost = input_cost + output_cost

        return {
            'input_tokens': self.total_input_tokens,
            'output_tokens': self.total_output_tokens,
            'input_cost_usd': input_cost,
            'output_cost_usd': output_cost,
            'total_cost_usd': total_cost
        }
